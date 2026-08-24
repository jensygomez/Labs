from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Crear Presentación 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Fundo claro (#F8FAFC)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)

# --- CABECERA ---
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "SYSTECH-HA-001 — Infrastructure Architecture"
p1.font.size = Pt(22)
p1.font.bold = True
p1.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)

p2 = tf.add_paragraph()
p2.text = "HA Web Cluster · NFSv4 + iSCSI Hybrid Storage · ZeroTier Mesh Overlay · Centralized DNS (dnsmasq) · Zabbix Monitoring"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

# --- PANEL IZQUIERDO: CONTROL DEVICE ---
ctrl_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.2), Inches(3.2), Inches(5.5))
ctrl_panel.fill.solid()
ctrl_panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
ctrl_panel.line.color.rgb = RGBColor(0xEB, 0xEE, 0xF2)
ctrl_panel.line.width = Pt(1.5)

# Badge Zabbix Control
badge_z = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(1.3), Inches(0.3), Inches(0.3))
badge_z.fill.solid()
badge_z.fill.fore_color.rgb = RGBColor(0xDC, 0x26, 0x26)
badge_z.line.fill.background()
pz = badge_z.text_frame.paragraphs[0]
pz.text = "Z"
pz.font.size = Pt(9)
pz.font.bold = True
pz.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
pz.alignment = PP_ALIGN.CENTER

# Texto Control Device
tb_ctrl = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3.0), Inches(0.6))
p = tb_ctrl.text_frame.paragraphs[0]
p.text = "CONTROL DEVICE\nAdministrator's Laptop"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
p.alignment = PP_ALIGN.CENTER

# Laptop Screen (Terminal)
screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.0), Inches(2.6), Inches(2.2))
screen.fill.solid()
screen.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
screen.line.fill.background()

tf_scr = screen.text_frame
tf_scr.margin_left = tf_scr.margin_top = Inches(0.15)
p_bash = tf_scr.paragraphs[0]
p_bash.text = "systech-control\n"
p_bash.font.size = Pt(10)
p_bash.font.bold = True
p_bash.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p_bash.alignment = PP_ALIGN.CENTER

p_cmd1 = tf_scr.add_paragraph()
p_cmd1.text = "$ ansible-playbook"
p_cmd1.font.size = Pt(9)
p_cmd1.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)

p_cmd2 = tf_scr.add_paragraph()
p_cmd2.text = "$ tofu apply"
p_cmd2.font.size = Pt(9)
p_cmd2.font.color.rgb = RGBColor(0x10, 0xB9, 0x81)

# Laptop Base
base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.2), Inches(1.0), Inches(0.1))
base.fill.solid()
base.fill.fore_color.rgb = RGBColor(0x94, 0xA3, 0xB8)
base.line.fill.background()

# Specs Terminal
tb_specs = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(3.0), Inches(1.5))
tf_sp = tb_specs.text_frame
p_sp = tf_sp.paragraphs[0]
p_sp.text = "Ansible Control Node\nOpenTofu + Ansible Vault\nPodman rootless container"
p_sp.font.size = Pt(9)
p_sp.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
p_sp.alignment = PP_ALIGN.CENTER


# --- PANEL DERECHO: PROXMOX VE ---
pve_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(1.2), Inches(9.1), Inches(5.5))
pve_panel.fill.solid()
pve_panel.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF6)
pve_panel.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
pve_panel.line.width = Pt(1.5)

# Header Proxmox
badge_p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(1.3), Inches(0.3), Inches(0.3))
badge_p.fill.solid()
badge_p.fill.fore_color.rgb = RGBColor(0xEA, 0x58, 0x0C)
badge_p.line.fill.background()
pp = badge_p.text_frame.paragraphs[0]
pp.text = "P"
pp.font.size = Pt(11)
pp.font.bold = True
pp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
pp.alignment = PP_ALIGN.CENTER

tb_pve = slide.shapes.add_textbox(Inches(4.35), Inches(1.25), Inches(3.0), Inches(0.4))
p = tb_pve.text_frame.paragraphs[0]
p.text = "Proxmox VE  Subnet 10.10.10.0/24"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

# Badge ZeroTier Proxmox
badge_z2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.4), Inches(1.3), Inches(0.3), Inches(0.3))
badge_z2.fill.solid()
badge_z2.fill.fore_color.rgb = RGBColor(0xDC, 0x26, 0x26)
badge_z2.line.fill.background()
pz2 = badge_z2.text_frame.paragraphs[0]
pz2.text = "Z"
pz2.font.size = Pt(9)
pz2.font.bold = True
pz2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
pz2.alignment = PP_ALIGN.CENTER


# --- COMPONENTES INTERNOS DE PROXMOX ---

# 1. CLIENT
client_node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(2.2), Inches(1.0), Inches(1.3))
client_node.fill.solid()
client_node.fill.fore_color.rgb = RGBColor(0x10, 0xB9, 0x81) # Emerald Green
client_node.line.fill.background()
tf_c = client_node.text_frame
tf_c.margin_top = Inches(0.3)
p = tf_c.paragraphs[0]
p.text = "client\n10.10.10.11\n\nTest traffic"
p.font.size = Pt(9)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 2. NUEVO NODO: DNS SERVER (dns01)
dns_node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(3.8), Inches(1.0), Inches(1.3))
dns_node.fill.solid()
dns_node.fill.fore_color.rgb = RGBColor(0x05, 0x96, 0x69) # Darker Emerald Green
dns_node.line.fill.background()
tf_dns = dns_node.text_frame
tf_dns.margin_top = Inches(0.15)
p = tf_dns.paragraphs[0]
p.text = "dns01\n10.10.10.20\n\ndnsmasq LXC\nLocal Storage"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 3. LOAD BALANCING LAYER (HA)
lb_layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(2.3), Inches(1.8))
lb_layer.fill.solid()
lb_layer.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEB)
lb_layer.line.color.rgb = RGBColor(0xFC, 0xE7, 0x88)

tb_lbl = slide.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(2.3), Inches(0.3))
p = tb_lbl.text_frame.paragraphs[0]
p.text = "LOAD BALANCING LAYER (HA)"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0xB4, 0x53, 0x09)
p.alignment = PP_ALIGN.CENTER

vip_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(2.15), Inches(1.7), Inches(0.4))
vip_box.fill.solid()
vip_box.fill.fore_color.rgb = RGBColor(0xFA, 0xCC, 0x15)
vip_box.line.fill.background()
p = vip_box.text_frame.paragraphs[0]
p.text = "Keepalived VIP\n10.10.10.30"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0x78, 0x35, 0x0F)
p.alignment = PP_ALIGN.CENTER

# LB Nodes
lb01 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(2.65), Inches(0.9), Inches(0.8))
lb01.fill.solid()
lb01.fill.fore_color.rgb = RGBColor(0xEA, 0x58, 0x0C)
lb01.line.fill.background()
p = lb01.text_frame.paragraphs[0]
p.text = "lb01\n10.10.10.21"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

lb02 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.45), Inches(2.65), Inches(0.9), Inches(0.8))
lb02.fill.solid()
lb02.fill.fore_color.rgb = RGBColor(0xEA, 0x58, 0x0C)
lb02.line.fill.background()
p = lb02.text_frame.paragraphs[0]
p.text = "lb02\n10.10.10.22"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 4. APPLICATION LAYER
app_layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(1.8), Inches(5.0), Inches(1.8))
app_layer.fill.solid()
app_layer.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
app_layer.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

tb_app = slide.shapes.add_textbox(Inches(7.7), Inches(1.85), Inches(5.0), Inches(0.3))
p = tb_app.text_frame.paragraphs[0]
p.text = "APPLICATION LAYER — Cluster (Apache + PHP)"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1D, 0x4E, 0xD8)
p.alignment = PP_ALIGN.CENTER

apps = [("app01", "10.10.10.31"), ("app02", "10.10.10.32"), ("app03", "10.10.10.33")]
for idx, (name, ip) in enumerate(apps):
    app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9 + (idx * 1.6)), Inches(2.2), Inches(1.4), Inches(1.2))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
    app_box.line.fill.background()
    tf_a = app_box.text_frame
    tf_a.margin_top = Inches(0.15)
    p = tf_a.paragraphs[0]
    p.text = f"{name}\n{ip}\n\nApache httpd\nPHP-pgsql\nNFSv4 client"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

# 5. MONITORING (Zabbix)
mon_layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.8), Inches(2.3), Inches(1.8))
mon_layer.fill.solid()
mon_layer.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
mon_layer.line.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

tb_mon = slide.shapes.add_textbox(Inches(5.2), Inches(3.85), Inches(2.3), Inches(0.3))
p = tb_mon.text_frame.paragraphs[0]
p.text = "MONITORING (Zabbix LXC)"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
p.alignment = PP_ALIGN.CENTER

zabbix_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(4.25), Inches(1.7), Inches(1.2))
zabbix_box.fill.solid()
zabbix_box.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
zabbix_box.line.fill.background()
tf_z = zabbix_box.text_frame
tf_z.margin_top = Inches(0.15)
p = tf_z.paragraphs[0]
p.text = "zabbix-lxc\n10.10.10.90\n\nPostgreSQL backend\nMonitoring all VMs"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 6. CENTRALIZED STORAGE
stg_layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.7), Inches(3.8), Inches(2.9), Inches(1.8))
stg_layer.fill.solid()
stg_layer.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
stg_layer.line.color.rgb = RGBColor(0xFC, 0xA5, 0xA5)

tb_stg = slide.shapes.add_textbox(Inches(7.7), Inches(3.85), Inches(2.9), Inches(0.3))
p = tb_stg.text_frame.paragraphs[0]
p.text = "CENTRALIZED STORAGE"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0x99, 0x1B, 0x1B)
p.alignment = PP_ALIGN.CENTER

stg_box = slide.shapes.add_shape(MSO_SHAPE.CAN, Inches(8.5), Inches(4.25), Inches(1.3), Inches(1.2))
stg_box.fill.solid()
stg_box.fill.fore_color.rgb = RGBColor(0xDC, 0x26, 0x26)
stg_box.line.fill.background()
tf_s = stg_box.text_frame
tf_s.margin_top = Inches(0.3)
p = tf_s.paragraphs[0]
p.text = "storage01\n10.10.10.50"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 7. DATABASE
db_layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(3.8), Inches(1.9), Inches(1.8))
db_layer.fill.solid()
db_layer.fill.fore_color.rgb = RGBColor(0xFA, 0xF5, 0xFF)
db_layer.line.color.rgb = RGBColor(0xE9, 0xD5, 0xFF)

tb_db = slide.shapes.add_textbox(Inches(10.8), Inches(3.85), Inches(1.9), Inches(0.3))
p = tb_db.text_frame.paragraphs[0]
p.text = "DATABASE"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = RGBColor(0x6B, 0x21, 0xA8)
p.alignment = PP_ALIGN.CENTER

db_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.1), Inches(4.25), Inches(1.3), Inches(1.2))
db_box.fill.solid()
db_box.fill.fore_color.rgb = RGBColor(0x93, 0x33, 0xEA)
db_box.line.fill.background()
tf_db = db_box.text_frame
tf_db.margin_top = Inches(0.3)
p = tf_db.paragraphs[0]
p.text = "db01\n10.10.10.40\n\nPostgreSQL 15+"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER


# --- LEYENDA (FOOTER) ---
leg_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.5))
leg_panel.fill.solid()
leg_panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
leg_panel.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

tb_leg = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5))
tf_leg = tb_leg.text_frame
p = tf_leg.paragraphs[0]
p.text = "Legend:   ● Client (traffic)   ● Centralized DNS (dns01)   ● HA Load Balancers   ● Application Cluster   ● Storage (NFS + iSCSI)   ● PostgreSQL   ● Zabbix Monitoring"
p.font.size = Pt(8)
p.font.color.rgb = RGBColor(0x47, 0x55, 0x69)

# Guardar Archivo
prs.save("SYSTECH_HA_001_Architecture_with_DNS.pptx")
